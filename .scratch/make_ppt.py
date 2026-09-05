# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x2E, 0x5E, 0xAA)
LBLUE = RGBColor(0xDD, 0xE9, 0xF7)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LGREEN = RGBColor(0xE1, 0xF2, 0xE1)
ORANGE = RGBColor(0xB4, 0x5F, 0x06)
LORANGE = RGBColor(0xFB, 0xE7, 0xCE)
GRAY = RGBColor(0x55, 0x55, 0x55)
BANNER = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if sub:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.4))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = sub
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = GRAY


def add_box(slide, x, y, w, h, title, lines, fill, line_color, font_pt=11):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(font_pt + 2)
    p0.font.bold = True
    p0.font.color.rgb = line_color
    p0.alignment = PP_ALIGN.LEFT
    for line in lines:
        p = tf.add_paragraph()
        p.text = "· " + line
        p.font.size = Pt(font_pt)
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.LEFT
    return shp


ARROW_XML = (
    '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{w}">'
    '<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
    '{dash}'
    '<a:tailEnd type="triangle" w="med" len="med"/>'
    '</a:ln>'
)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width_pt=1.75, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    hexcolor = "%02X%02X%02X" % (color[0], color[1], color[2])
    dash_xml = '<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dash"/>' if dashed else ""
    xml = ARROW_XML.format(w=int(width_pt * 12700), color=hexcolor, dash=dash_xml)
    new_ln = parse_xml(xml)
    spPr = conn._element.spPr
    old_ln = spPr.find(qn('a:ln'))
    if old_ln is not None:
        spPr.remove(old_ln)
    spPr.append(new_ln)
    return conn


def add_label(slide, x, y, w, h, text, size=9.5, color=GRAY, bold=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


# ================================================================= Slide 1
s1 = prs.slides.add_slide(BLANK)
add_title(s1, "1. 턴 안에서 쓰는 정보 (in-turn)",
          "질문 1건에 답하는 동안만 쓰이는 흐름 — 이 자체는 아무것도 새로 저장하지 않는다")

cx = 6.67          # 중앙 세로줄 x좌표(인치)
w1 = 4.6
x1 = cx - w1 / 2

y = 1.30
add_box(s1, x1, y, w1, 0.55, "사용자 질문", [], LBLUE, BLUE, font_pt=13)
add_arrow(s1, cx, y + 0.55, cx, y + 0.80)

y = 2.10
add_box(s1, x1 - 0.5, y, w1 + 1.0, 1.05, "[항상 자동] _prefetch_memory()", [
    "memory.search(tier=\"all\") — STM+MTM+LTM 조회수 무관 즉시 조회",
    "시스템 프롬프트 뒤 user 메시지로 강제 주입",
    "도구 아님 — LLM이 부를지 판단하는 대상 아님",
], LORANGE, ORANGE, font_pt=10.5)
add_arrow(s1, cx, y + 1.05, cx, y + 1.30)

y = 3.40
add_box(s1, x1 - 0.9, y, w1 + 1.8, 0.55, "LLM 추론 루프 (max_tool_calls 회까지)", [], LBLUE, BLUE, font_pt=12)
loop_bottom = y + 0.55

# 두 도구로 분기
tool_y = 4.30
left_x, left_w = 0.7, 3.1
right_x, right_w = 9.55, 3.1
add_arrow(s1, cx - 0.9, loop_bottom, left_x + left_w / 2, tool_y)
add_arrow(s1, cx + 0.9, loop_bottom, right_x + right_w / 2, tool_y)
add_label(s1, x1 - 0.9, loop_bottom + 0.03, w1 + 1.8, 0.28, "tool_calls 있음", size=9, color=GRAY)

add_box(s1, left_x, tool_y, left_w, 1.05, "search_documents", [
    "doc_rag: Dense+BM25", "→ RRF (리랭커 없음)", "→ 문서 청크 근거카드",
], LGREEN, GREEN, font_pt=10)

add_box(s1, right_x, tool_y, right_w, 1.05, "search_graph", [
    "graph_rag: 닫힌 어휘 개체 해소", "(정확매칭 → 퍼지 fallback)", "→ N홉 → 근거카드",
], LGREEN, GREEN, font_pt=10)

# 도구 -> 루프로 복귀 (다시 LLM이 판단)
back_y = 5.55
add_arrow(s1, left_x + left_w / 2, tool_y + 1.05, cx - 0.9, back_y)
add_arrow(s1, right_x + right_w / 2, tool_y + 1.05, cx + 0.9, back_y)
add_label(s1, x1 - 0.9, back_y + 0.02, w1 + 1.8, 0.26, "tool_calls 없음 → 최종 답변", size=9, color=GRAY)

y = 5.85
add_box(s1, x1 - 0.9, y, w1 + 1.8, 0.55, "최종 답변 생성", [], LBLUE, BLUE, font_pt=13)
add_arrow(s1, cx, y + 0.55, cx, y + 0.80)

y = 6.65
add_box(s1, x1 - 1.5, y, w1 + 3.0, 0.65, "_score_attribution()", [
    "답변에 쓰인 근거 비중 계산 → 메모리 카드만 memory.credit() → 조회수(views) ↑",
], LORANGE, ORANGE, font_pt=10)

# ================================================================= Slide 2
s2 = prs.slides.add_slide(BLANK)
add_title(s2, "2. 메모리 / RAG에 입력되는 개념",
          "경로 A(대화→STM, 매 턴 자동·승격 구조)  vs  경로 B(문서→RAG, 배치 실행·승격 없음)")

ax, aw = 0.55, 5.95
bx, bw = 6.85, 5.95

add_label(s2, ax, 1.15, aw, 0.35, "경로 A: 대화(chat) → STM", size=15, color=BLUE, bold=True, align=PP_ALIGN.LEFT)
add_label(s2, bx, 1.15, bw, 0.35, "경로 B: 문서 코퍼스 → doc_rag / graph_rag", size=15, color=ORANGE, bold=True, align=PP_ALIGN.LEFT)

# ---- 경로 A
y = 1.60
add_box(s2, ax, y, aw, 0.55, "turn 종료 (question, answer)", [], LBLUE, BLUE, font_pt=11)
add_arrow(s2, ax + aw / 2, y + 0.55, ax + aw / 2, y + 0.75)

y = 2.35
add_box(s2, ax, y, aw, 0.95, "SessionRecorder.record_turn()", [
    "data/raw/ai_chat/<session_id>.jsonl 원문 append",
    "FactExtractor.extract() — 이번 턴만 보고 원자 사실 추출",
], LGREEN, GREEN, font_pt=10)
add_arrow(s2, ax + aw / 2, y + 0.95, ax + aw / 2, y + 1.15)

y = 3.50
add_box(s2, ax, y, aw, 0.55, "memory.add(tier=\"stm\") → facts.sqlite3", [], LBLUE, BLUE, font_pt=11)
add_arrow(s2, ax + aw / 2, y + 0.55, ax + aw / 2, y + 0.75)

y = 4.25
add_box(s2, ax, y, aw, 0.75, "_check_promotion()", [
    "promote_touched() — 이번 턴 조회수 상승분만 확인 → 조건 충족 시 STM→MTM 즉시 승격",
], LORANGE, ORANGE, font_pt=10)
add_arrow(s2, ax + aw / 2, y + 0.75, ax + aw / 2, y + 0.95, dashed=True)
add_label(s2, ax, y + 0.97, aw, 0.25, "별도 배치 실행 (memory_promote.run() / --ltm)", size=9)

y = 5.45
add_box(s2, ax, y, aw, 0.95, "배치 승격 · 정리", [
    "STM 폐기(창 밖+조회수 미달) · MTM 중복 통합",
    "MTM→LTM: 3기준 통과 + 기존 LTM 대조 → 신규/중복폐기/갱신",
], LGREEN, GREEN, font_pt=10)

# ---- 경로 B
y = 1.60
add_box(s2, bx, y, bw, 0.55, "config/doc_rag.yaml corpus.roots", [
    "프로젝트 파일 폴더 — 채팅과 무관",
], LBLUE, BLUE, font_pt=10.5)
add_arrow(s2, bx + bw / 2, y + 0.55, bx + bw / 2, y + 0.75)

y = 2.35
add_box(s2, bx, y, bw, 1.05, "python -m doc_rag.ingest  (오프라인 배치)", [
    "walk_corpus → Docling 변환 → 형식별 청킹",
    "임베딩(mE5-base)",
    "→ index/doc_rag/{chunks.sqlite3, embeddings.npy}",
], LGREEN, GREEN, font_pt=10)
add_arrow(s2, bx + bw / 2, y + 1.05, bx + bw / 2, y + 1.25)

y = 3.60
add_box(s2, bx, y, bw, 0.95, "python -m graph_rag.build  (오프라인 배치)", [
    "doc_rag 청크에서 닫힌 어휘(schema.py) 개체 표기 탐지",
    "구조 엣지 + MENTIONS + CO_OCCURS → graph.sqlite3",
], LGREEN, GREEN, font_pt=10)
add_arrow(s2, bx + bw / 2, y + 0.95, bx + bw / 2, y + 1.15)

y = 4.75
add_box(s2, bx, y, bw, 0.80, "= LTM급 근거로 즉시 취급", [
    "승격 절차 없음 — doc_rag 자체가 에이전트의 장기기억(LTM) 역할",
], LORANGE, ORANGE, font_pt=10.5)

# 하단 배너
add_box(s2, 0.55, 6.60, 12.25, 0.65, "핵심 차이", [
    "A는 매 턴 자동 실행 + STM→MTM→LTM 승격(살아있는 계층) / B는 CLI로 명시 실행해야 갱신되는 정적 인덱스, 승격 없음",
], BANNER, NAVY, font_pt=10.5)

out = r"D:\LAB_RAG\Org-AI-Body\.scratch\info_flow.pptx"
prs.save(out)
print("saved:", out)
