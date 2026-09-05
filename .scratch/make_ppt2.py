# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x2E, 0x5E, 0xAA)
LBLUE = RGBColor(0xDD, 0xE9, 0xF7)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LGREEN = RGBColor(0xE1, 0xF2, 0xE1)
ORANGE = RGBColor(0xB4, 0x5F, 0x06)
LORANGE = RGBColor(0xFB, 0xE7, 0xCE)
RED = RGBColor(0xA6, 0x2C, 0x2C)
LRED = RGBColor(0xFA, 0xDE, 0xDE)
PURPLE = RGBColor(0x6A, 0x3D, 0x9A)
LPURPLE = RGBColor(0xEC, 0xE1, 0xF6)
GRAY = RGBColor(0x55, 0x55, 0x55)
BANNER = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if sub:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.68), Inches(12.3), Inches(0.4))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = sub
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY


def add_box(slide, x, y, w, h, title, lines, fill, line_color, font_pt=10.5, title_pt=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.11)
    tf.margin_right = Inches(0.11)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(title_pt or (font_pt + 1.5))
    p0.font.bold = True
    p0.font.color.rgb = line_color
    p0.alignment = PP_ALIGN.LEFT
    for line in lines:
        p = tf.add_paragraph()
        p.text = "· " + line
        p.font.size = Pt(font_pt)
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.LEFT
    return shp


ARROW_XML = (
    '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{w}" cap="rnd">'
    '<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
    '{dash}'
    '<a:round/>'
    '<a:tailEnd type="triangle" w="lg" len="lg"/>'
    '</a:ln>'
)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width_pt=2.5, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    hexcolor = "%02X%02X%02X" % (color[0], color[1], color[2])
    dash_xml = '<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dash"/>' if dashed else ""
    xml = ARROW_XML.format(w=int(width_pt * 12700), color=hexcolor, dash=dash_xml)
    new_ln = parse_xml(xml)
    spPr = conn._element.spPr
    old_ln = spPr.find(qn('a:ln'))
    if old_ln is not None:
        spPr.remove(old_ln)
    spPr.append(new_ln)
    return conn


def add_label(slide, x, y, w, h, text, size=9.5, color=GRAY, bold=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
    return box


def footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = GRAY


# ================================================================= Slide 0: 표지 / 목차
s0 = prs.slides.add_slide(BLANK)
add_title(s0, "설계 근거 정리", "Org-AI-Body 메모리·검색 파이프라인 — 5개 주제")
items = [
    ("1", "MTM → LTM 승격 방식 정의", BLUE, LBLUE),
    ("2", "문서 코퍼스 입력 시 처리 분기 정의", GREEN, LGREEN),
    ("3", "지식그래프 검색(엔티티 해소) 성능 개선", ORANGE, LORANGE),
    ("4", "임베딩/리랭커 모델 메모리 탑재·재사용", PURPLE, LPURPLE),
    ("5", "STM/MTM/LTM 분기 구조 — 타 아키텍처 비교", RED, LRED),
]
y = 1.55
for num, text, c, lc in items:
    circ = s0.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y), Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = c
    circ.line.fill.background()
    tf = circ.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    bar = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.75), Inches(y + 0.03), Inches(10.6), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = lc
    bar.line.color.rgb = c
    bar.line.width = Pt(1)
    bar.shadow.inherit = False
    tf2 = bar.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf2.margin_left = Inches(0.2)
    p2 = tf2.paragraphs[0]
    p2.text = text
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    y += 0.98

# ================================================================= Slide 1: MTM -> LTM
s1 = prs.slides.add_slide(BLANK)
add_title(s1, "1. MTM → LTM 승격 방식",
          "3개 기준(AND) 통과 + 기존 LTM과 1:1 대조 — 흡수 병합이 아니라 신규/중복폐기/갱신 판정")

add_box(s1, 0.5, 1.35, 3.7, 1.5, "MTM 사실", [
    "출처(traceable) — source_id가 원문 로그로",
    "  이어지는가 (데이터로 확인, 결정적)",
    "진실성(truthful) — LLM 판정 (문서 근거 힌트)",
    "유용성(useful) — LLM 판정",
], LBLUE, BLUE, font_pt=10)
add_label(s1, 0.5, 2.9, 3.7, 0.3, "셋 다 True → 통과 / 하나라도 False → MTM 잔류", size=9)

add_arrow(s1, 4.2, 2.05, 4.85, 2.05)
add_box(s1, 4.9, 1.55, 2.9, 1.0, "승격 후보", [
    "3기준 통과한 MTM 사실",
], LGREEN, GREEN, font_pt=10.5)

add_arrow(s1, 6.35, 2.55, 6.35, 2.9)
add_box(s1, 3.4, 2.9, 5.9, 0.55, "기존 LTM 전체와 1:1 대조 (_resolve_against_ltm)", [], LPURPLE, PURPLE, font_pt=11)

by = 3.75
bw = 2.75
gap = 0.15
bx0 = 0.5
branches = [
    ("무관한 사안", ["코사인 < 0.65"], "→ 그대로 신규 승격", LGREEN, GREEN),
    ("같은 사실 (중복)", ["코사인 >= 0.65", "상호함의 >= 2.0"], "→ 후보 폐기\n(기존 LTM 안 건드림)", LRED, RED),
    ("같은 사안·다른 정보", ["코사인 >= 0.65", "상호함의 [-6.0, 2.0)"], "→ 기존 LTM 레코드 갱신", LORANGE, ORANGE),
    ("표면만 겹친 무관 문장", ["코사인 >= 0.65", "상호함의 < -6.0"], "→ 신규 승격", LBLUE, BLUE),
]
for i, (title, conds, outcome, fill, line) in enumerate(branches):
    bx = bx0 + i * (bw + gap)
    add_arrow(s1, 6.35, 3.45, bx + bw / 2, by)
    add_box(s1, bx, by, bw, 1.15, title, conds, fill, line, font_pt=10.5, title_pt=11.5)
    add_label(s1, bx, by + 1.18, bw, 0.55, outcome, size=10, color=NAVY, bold=True)

footer(s1, "UPDATE_FLOOR=-6.0은 실측 보정: 코사인 0.85·상호함의 -7.65(완전 무관, \"담당자 변경\" vs \"전담기관명\")로 잘못 갱신되던 사례를 발견하고 도입 · memory_promote.py:267")

# ================================================================= Slide 2: 문서 입력 분기
s2 = prs.slides.add_slide(BLANK)
add_title(s2, "2. 문서 코퍼스 입력 시 처리 분기",
          "4중 스킵 필터 → 3단 변환 폴백 → 형식별 청킹")

y = 1.35
add_box(s2, 0.5, y, 12.3, 0.55, "파일 1건", [], LBLUE, BLUE, font_pt=12)
add_arrow(s2, 6.65, y + 0.55, 6.65, y + 0.75)

y = 2.10
skip_w = 2.95
skip_labels = [
    ("확장자 미지원", "skip(unsupported)"),
    ("크기 > 120MB", "skip(too_large)"),
    ("캐시 일치(mtime 동일)", "skip(cached)"),
    ("SHA256 중복 문서", "skip(duplicate_of)"),
]
for i, (cond, out) in enumerate(skip_labels):
    bx = 0.5 + i * (skip_w + 0.13)
    add_box(s2, bx, y, skip_w, 0.75, cond, [out], LRED, RED, font_pt=9.5, title_pt=10)
add_label(s2, 0.5, y + 0.78, 12.3, 0.25, "4개 필터 모두 통과한 파일만 아래로 진행", size=9.5, bold=True)
add_arrow(s2, 6.65, y + 0.75 + 0.25, 6.65, 3.35)

w_conv = 3.8
cx0 = 6.65 - w_conv / 2

y = 3.35
add_box(s2, cx0, y, w_conv, 0.65, "1차: docling-parse 변환", [], LGREEN, GREEN, font_pt=11)
add_arrow(s2, 6.65, y + 0.65, 6.65, y + 1.02)
add_label(s2, cx0 + w_conv + 0.15, y + 0.65, 3.0, 0.4, "실패 시 (PDF만)  ▶", size=10, color=RED, bold=True, align=PP_ALIGN.LEFT)

y = 4.35
add_box(s2, cx0, y, w_conv, 0.65, "2차: pypdfium 백엔드 재시도", [], LORANGE, ORANGE, font_pt=11)
add_arrow(s2, 6.65, y + 0.65, 6.65, y + 1.02)
add_label(s2, cx0 + w_conv + 0.15, y + 0.65, 3.4, 0.4, "그래도 청크 0개(스캔 PDF)  ▶", size=10, color=RED, bold=True, align=PP_ALIGN.LEFT)

y = 5.35
add_box(s2, cx0, y, w_conv, 0.65, "3차: EasyOCR 변환 (이때만 OCR 비용)", [], LPURPLE, PURPLE, font_pt=11)
add_arrow(s2, 6.65, y + 0.65, 6.65, y + 1.00)

y = 6.35
add_box(s2, 6.65 - w_conv / 2 - 1.3, y, w_conv + 2.6, 0.68,
        "형식별 청킹 (by_format)",
        ["pptx 320/32(짧게) · hwp 640/80(길게) · xlsx 768/0(overlap 없음, config만 확인)"],
        LBLUE, BLUE, font_pt=10)

footer(s2, "실측: PDF 552개 중 306개가 텍스트 레이어 없는 스캔본 → 처음부터 전체 OCR은 낭비 · doc_rag/ingest.py:148")

# ================================================================= Slide 3: 그래프 검색 개선
s3 = prs.slides.add_slide(BLANK)
add_title(s3, "3. 지식그래프 검색(엔티티 해소) 성능 개선",
          "정확일치 우선 → 실패 시에만 문자 단위 퍼지 매칭 폴백")

add_box(s3, 0.6, 1.4, 4.3, 1.0, "질의문", ["예: \"경북대산학협력단 담당자는?\"" ], LBLUE, BLUE, font_pt=10.5)
add_arrow(s3, 4.9, 1.9, 5.5, 1.9)

add_box(s3, 5.55, 1.4, 3.4, 1.0, "표기 완전일치 시도", [
    "긴 표기부터, 빠르고 오탐 없음",
], LGREEN, GREEN, font_pt=10)
add_arrow(s3, 7.25, 2.4, 7.25, 2.75)
add_label(s3, 5.55, 2.4, 3.4, 0.3, "못 찾음", size=9.5, bold=True)

add_box(s3, 3.0, 2.75, 8.4, 1.55, "퍼지 매칭 폴백 (완전일치 실패했을 때만)", [
    "개체 표기와 길이가 비슷한 질의문 부분 문자열 전부 비교 (difflib.SequenceMatcher)",
    "개체별 최고 점수 >= 0.65 → 채택",
], LORANGE, ORANGE, font_pt=10.5)

add_arrow(s3, 7.25, 4.30, 7.25, 4.6)
add_box(s3, 1.0, 4.6, 5.8, 1.65, "왜 임베딩(의미 유사도)이 아니라 문자 매칭인가", [
    "질의 임베딩 vs 개체 centroid 시도했으나, 단일 프로젝트",
    "코퍼스라 모든 centroid가 뭉쳐 변별력 없음",
    "(정답·오답 코사인 차이 0.001~0.1) → 폐기",
], LPURPLE, PURPLE, font_pt=10.5)

add_box(s3, 7.0, 4.6, 5.8, 1.65, "왜 0.65인가 (실측 근거)", [
    "개체 실제 언급된 질의: 0.800~1.000",
    "언급 없는 질의: 0.333~0.400",
    "0.6은 동률 구간(오답 기관들도 같이 걸림) → 0.65로 회피",
], LGREEN, GREEN, font_pt=10.5)

footer(s3, "graph_rag/query.py:28 FUZZY_MATCH_THRESHOLD=0.65 — 완전일치·퍼지매칭을 항상 같이 돌리지 않는 이유: 완전일치 결과가 있으면 표기가 정확해 퍼지 매칭이 낄 자리가 없음(중복 계산 방지)")

# ================================================================= Slide 4: 임베더 재사용
s4 = prs.slides.add_slide(BLANK)
add_title(s4, "4. 임베딩/리랭커 모델 메모리 탑재·재사용",
          "STM·MTM·LTM·그래프·문서 검색이 전부 같은 임베딩 공간을 쓰도록 모델 하나로 통일")

callers = [
    ("FactMemory", "STM/MTM/LTM 사실 검색·승급 판정", LBLUE, BLUE),
    ("DocRetriever", "문서 청크 검색(doc_rag)", LGREEN, GREEN),
    ("GraphRetriever.query", "엔티티 퍼지 매칭 근거 재정렬", LORANGE, ORANGE),
    ("GraphRetriever.build", "그래프 구축 시 청크 임베딩", LPURPLE, PURPLE),
    ("AgentRuntime", "답변 기여도(attribution) 계산", LRED, RED),
]
cy = 1.55
ch = 0.72
target_x, target_y0, target_h = 5.65, 3.2, 1.5
for i, (name, desc, fill, line) in enumerate(callers):
    y = cy + i * (ch + 0.12)
    add_box(s4, 0.5, y, 4.3, ch, name, [desc], fill, line, font_pt=10, title_pt=11.5)
    target_y = target_y0 + target_h * (i + 0.5) / len(callers)
    add_arrow(s4, 4.85, y + ch / 2, target_x, target_y)

add_box(s4, 5.7, 3.2, 4.6, 1.5, "embedding_factory.build_embedder(cfg)", [
    "캐시 키 = (backend, model_id, device, batch_size, max_seq_length)",
    "같은 key → 이미 로드된 인스턴스 그대로 반환",
], LBLUE, BLUE, font_pt=10)

add_arrow(s4, 8.0, 4.7, 8.0, 5.05)
add_box(s4, 5.7, 5.05, 4.6, 0.85, "모델 하나로 통일된 임베딩 공간", [
    "다섯 저장소가 같은 벡터 공간 → 코사인 비교(중복판정·승급·퍼지매칭)가 일관됨",
], LGREEN, GREEN, font_pt=9.5)

add_box(s4, 10.55, 3.2, 2.3, 2.7, "실측 낭비 (캐시 없을 때)", [
    "VRAM 8.47GB 사용",
    "그중 4.24GB 중복",
    "로딩 7.2초씩 중복",
    "(arctic-ko 3벌 중복 로드)",
], LRED, RED, font_pt=9.5)

footer(s4, "인스턴스는 무상태(모델 가중치 + 디스크 캐시 핸들뿐) — 스레드 간 공유해도 문제 없음 · doc_rag/embedding_factory.py")

# ================================================================= Slide 5: 계층 분기 비교
s5 = prs.slides.add_slide(BLANK)
add_title(s5, "5. STM/MTM/LTM 분기 구조 — 타 아키텍처 비교",
          "현재: 매 턴 세 계층 전부 조회(예산 3/3/3, LLM 판단 아님) — 다른 구조는 어떻게 분기하는가")

rows = [
    ("Self-RAG", "모델 자체(파인튜닝)", "[Retrieve]/[IsSupported]/[IsUseful] 특수 토큰 직접 생성", "불가 — API 전용 모델 파인튜닝 불가", RED),
    ("FLARE", "생성 중 신뢰도", "토큰 생성 확률이 임계치 이하로 떨어질 때만 검색 트리거", "불가 — API는 토큰별 로그프롭 미노출", RED),
    ("Adaptive-RAG", "별도 학습된 분류기", "질의 복잡도로 검색없음/1단계/다단계 라우팅", "가능 — 라벨 데이터 필요(비용)", ORANGE),
    ("Probing-RAG", "은닉 상태 프로브", "은닉층 활성값으로 검색 필요 예측, 57.5% 스킵", "불가 — 은닉 상태 API 접근 불가", RED),
    ("Letta/MemGPT", "토큰 압박 신호", "컨텍스트 70% 초과 시 시스템이 경고 주입 → LLM 함수호출", "부분 가능 — 판단 시점 게이트 아이디어만", ORANGE),
    ("Generative Agents", "분기 없음", "recency·importance·relevance 단일 점수로 전체 랭킹", "참고용 — 계층 개념 자체가 없음", GRAY),
]
hy = 1.35
header_shapes = [
    add_box(s5, 0.5, hy, 2.3, 0.5, "시스템", [], NAVY, NAVY, font_pt=10, title_pt=11),
    add_box(s5, 2.85, hy, 2.6, 0.5, "판단 주체", [], NAVY, NAVY, font_pt=10, title_pt=11),
    add_box(s5, 5.5, hy, 4.4, 0.5, "판단 근거", [], NAVY, NAVY, font_pt=10, title_pt=11),
    add_box(s5, 10.0, hy, 2.8, 0.5, "우리 스택 적용", [], NAVY, NAVY, font_pt=10, title_pt=11),
]
for shp in header_shapes:
    shp.fill.fore_color.rgb = NAVY
    shp.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

ry = hy + 0.55
rh = 0.72
for name, who, why, ours, c in rows:
    add_box(s5, 0.5, ry, 2.3, rh, name, [], BANNER, c, font_pt=10.5, title_pt=11)
    add_box(s5, 2.85, ry, 2.6, rh, who, [], BANNER, GRAY, font_pt=10, title_pt=10)
    add_box(s5, 5.5, ry, 4.4, rh, why, [], BANNER, GRAY, font_pt=9.5, title_pt=9.5)
    add_box(s5, 10.0, ry, 2.8, rh, ours, [], BANNER, c, font_pt=9.5, title_pt=9.5)
    ry += rh + 0.05

add_box(s5, 0.5, ry + 0.08, 12.3, 0.7,
        "결론 — 즉시 적용 가능한 방향",
        ["모델 내부 접근 필요한 방식은 전부 배제 → 이미 계산한 코사인 점수에 하한선(예: 0.5)을 두고, 그 밑이면 해당 계층 결과를 프롬프트에서 제외 (추가 비용 0)"],
        LGREEN, GREEN, font_pt=10.5)

out = r"D:\LAB_RAG\Org-AI-Body\.scratch\design_review.pptx"
prs.save(out)
print("saved:", out)
