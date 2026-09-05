# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x2E, 0x5E, 0xAA)
LBLUE = RGBColor(0xDD, 0xE9, 0xF7)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LGREEN = RGBColor(0xE1, 0xF2, 0xE1)
ORANGE = RGBColor(0xB4, 0x5F, 0x06)
LORANGE = RGBColor(0xFB, 0xE7, 0xCE)
RED = RGBColor(0xA6, 0x2C, 0x2C)
LRED = RGBColor(0xFA, 0xDE, 0xDE)
PURPLE = RGBColor(0x6A, 0x3D, 0x9A)
LPURPLE = RGBColor(0xEC, 0xE1, 0xF6)
GRAY = RGBColor(0x55, 0x55, 0x55)
BANNER = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---- 최소 폰트 하한 (이 밑으로는 절대 안 내림) ----
MIN_BODY = 13
MIN_TITLE = 15


def add_title(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(27)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if sub:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.4))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = sub
        p2.font.size = Pt(13.5)
        p2.font.color.rgb = GRAY


def add_box(slide, x, y, w, h, title, lines, fill, line_color, font_pt=MIN_BODY, title_pt=None,
            align=PP_ALIGN.LEFT):
    font_pt = max(font_pt, MIN_BODY)
    title_pt = max(title_pt or (font_pt + 2), MIN_TITLE)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(2)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(title_pt)
    p0.font.bold = True
    p0.font.color.rgb = line_color
    p0.alignment = align
    for line in lines:
        p = tf.add_paragraph()
        p.text = "•  " + line
        p.font.size = Pt(font_pt)
        p.font.color.rgb = NAVY
        p.alignment = align
        p.space_before = Pt(3)
    return shp


ARROW_XML = (
    '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{w}" cap="rnd">'
    '<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
    '{dash}'
    '<a:round/>'
    '<a:tailEnd type="triangle" w="lg" len="lg"/>'
    '</a:ln>'
)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width_pt=3.0, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    hexcolor = "%02X%02X%02X" % (color[0], color[1], color[2])
    dash_xml = '<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dash"/>' if dashed else ""
    xml = ARROW_XML.format(w=int(width_pt * 12700), color=hexcolor, dash=dash_xml)
    new_ln = parse_xml(xml)
    spPr = conn._element.spPr
    old_ln = spPr.find(qn('a:ln'))
    if old_ln is not None:
        spPr.remove(old_ln)
    spPr.append(new_ln)
    return conn


def add_label(slide, x, y, w, h, text, size=MIN_BODY, color=GRAY, bold=False, align=PP_ALIGN.CENTER):
    size = max(size, 11)  # 라벨은 보조 텍스트라 11pt까지만 허용
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
    return box


def footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.14), Inches(12.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = GRAY


# ================================================================= Slide 0: 표지
s0 = prs.slides.add_slide(BLANK)
add_title(s0, "설계 근거 정리", "Org-AI-Body 메모리·검색 파이프라인 — 5개 주제")
items = [
    ("1", "MTM → LTM 승격 방식 정의", BLUE, LBLUE),
    ("2", "문서 코퍼스 입력 시 처리 분기 정의", GREEN, LGREEN),
    ("3", "지식그래프 검색(엔티티 해소) 성능 개선", ORANGE, LORANGE),
    ("4", "임베딩/리랭커 모델 메모리 탑재·재사용", PURPLE, LPURPLE),
    ("5", "STM/MTM/LTM 분기 구조 — 타 아키텍처 비교", RED, LRED),
]
y = 1.65
for num, text, c, lc in items:
    circ = s0.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y), Inches(0.68), Inches(0.68))
    circ.fill.solid()
    circ.fill.fore_color.rgb = c
    circ.line.fill.background()
    tf = circ.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    bar = s0.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.85), Inches(y), Inches(10.55), Inches(0.68))
    bar.fill.solid()
    bar.fill.fore_color.rgb = lc
    bar.line.color.rgb = c
    bar.line.width = Pt(1.5)
    bar.shadow.inherit = False
    tf2 = bar.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf2.margin_left = Inches(0.25)
    p2 = tf2.paragraphs[0]
    p2.text = text
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    y += 1.02

# ================================================================= Slide 1: MTM -> LTM
s1 = prs.slides.add_slide(BLANK)
add_title(s1, "1. MTM → LTM 승격 방식",
          "3개 기준(AND) 통과 → 기존 LTM과 1:1 대조해 신규/중복폐기/갱신 판정")

add_box(s1, 0.5, 1.35, 12.3, 1.5, "1단계 — 3개 기준 AND 게이트", [
    "출처(traceable) — source_id → 원문 로그 매칭 (결정적)",
    "진실성·유용성(truthful/useful) — LLM 판정 (문서 근거 힌트)",
], LBLUE, BLUE, font_pt=14, title_pt=16)

add_arrow(s1, 6.65, 2.85, 6.65, 3.15)
add_box(s1, 3.9, 3.15, 5.5, 0.65, "2단계 — 기존 LTM 1:1 대조", [], LPURPLE, PURPLE, font_pt=14, title_pt=15)
add_arrow(s1, 6.65, 3.80, 6.65, 4.10)

by = 4.10
bw, bh = 5.95, 1.42
gapx, gapy = 0.4, 0.25
branches = [
    ("무관한 사안", "주제 유사도 낮음", "→ 신규 승격", LGREEN, GREEN),
    ("같은 사실(중복)", "주제 유사 · 내용 일치", "→ 후보 폐기 (기존 LTM 유지)", LRED, RED),
    ("같은 사안·다른 정보", "주제 유사 · 내용 다름", "→ 기존 LTM 갱신", LORANGE, ORANGE),
    ("표면만 겹친 무관", "주제 유사 · 내용 무관", "→ 신규 승격", LBLUE, BLUE),
]
for i, (title, cond, outcome, fill, line) in enumerate(branches):
    col, row = i % 2, i // 2
    bx = 0.5 + col * (bw + gapx)
    yy = by + row * (bh + gapy)
    add_box(s1, bx, yy, bw, bh, title, [cond, outcome], fill, line, font_pt=13, title_pt=15.5)

footer(s1, "UPDATE_FLOOR=-6.0 = 실측 보정치 (오분류 사례: 코사인 0.85·상호함의 -7.65, \"담당자 변경\" vs \"전담기관명\") · memory_promote.py:267")

# ================================================================= Slide 2: 문서 입력 분기
s2 = prs.slides.add_slide(BLANK)
add_title(s2, "2. 문서 코퍼스 입력 시 처리 분기",
          "4중 스킵 필터 → 3단 변환 폴백 → 형식별 청킹")

add_box(s2, 4.4, 1.3, 4.5, 0.6, "파일 1건", [], LBLUE, BLUE, font_pt=14, title_pt=16, align=PP_ALIGN.CENTER)
add_arrow(s2, 6.65, 1.9, 6.65, 2.15)

skip_w, skip_h = 5.95, 0.85
skip_labels = [
    ("확장자 미지원 → skip", ORANGE),
    ("크기 > 120MB → skip", ORANGE),
    ("캐시 일치(mtime 동일) → skip", GRAY),
    ("SHA256 중복 문서 → skip", GRAY),
]
for i, (text, c) in enumerate(skip_labels):
    col, row = i % 2, i // 2
    bx = 0.5 + col * (skip_w + 0.4)
    yy = 2.2 + row * (skip_h + 0.15)
    add_box(s2, bx, yy, skip_w, skip_h, text, [], LRED, c, font_pt=13, title_pt=14, align=PP_ALIGN.CENTER)

add_arrow(s2, 6.65, 4.05, 6.65, 4.30)
add_label(s2, 0.5, 4.06, 12.3, 0.3, "4개 필터 통과 → 진행", size=12, bold=True)

w_conv = 6.4
cx0 = 6.65 - w_conv / 2
y = 4.35
add_box(s2, cx0, y, w_conv, 0.62, "1차: docling-parse 변환", [], LGREEN, GREEN, font_pt=13.5, title_pt=14.5, align=PP_ALIGN.CENTER)
add_label(s2, cx0 + w_conv + 0.15, y - 0.02, 3.0, 0.7, "실패 시\n(PDF만)", size=11.5, color=RED, bold=True, align=PP_ALIGN.LEFT)
add_arrow(s2, 6.65, y + 0.62, 6.65, y + 0.95)

y = 5.10
add_box(s2, cx0, y, w_conv, 0.62, "2차: pypdfium 백엔드 재시도", [], LORANGE, ORANGE, font_pt=13.5, title_pt=14.5, align=PP_ALIGN.CENTER)
add_label(s2, cx0 + w_conv + 0.15, y - 0.02, 3.0, 0.7, "청크 0개\n(스캔 PDF)", size=11.5, color=RED, bold=True, align=PP_ALIGN.LEFT)
add_arrow(s2, 6.65, y + 0.62, 6.65, y + 0.95)

y = 5.85
add_box(s2, cx0, y, w_conv, 0.62, "3차: EasyOCR (OCR 비용 이 단계만)", [], LPURPLE, PURPLE, font_pt=13, title_pt=14, align=PP_ALIGN.CENTER)

footer(s2, "실측: PDF 552개 중 306개 텍스트 레이어 없는 스캔본 → 실패 시에만 다음 단계 전환(지연 생성) · doc_rag/ingest.py:148")

# ================================================================= Slide 3: 그래프 검색 개선
s3 = prs.slides.add_slide(BLANK)
add_title(s3, "3. 지식그래프 검색(엔티티 해소) 성능 개선",
          "정확일치 우선 → 실패 시에만 문자 단위 퍼지 매칭 폴백")

add_box(s3, 0.5, 1.35, 5.9, 0.9, "질의문", ["예: \"경북대산학협력단 담당자는?\""], LBLUE, BLUE, font_pt=13, title_pt=15)
add_arrow(s3, 6.55, 1.8, 6.9, 1.8)
add_box(s3, 6.95, 1.35, 5.9, 0.9, "표기 완전일치 시도", ["긴 표기 우선 · 빠름 · 오탐 없음"], LGREEN, GREEN, font_pt=13, title_pt=15)

add_arrow(s3, 9.9, 2.25, 9.9, 2.55)
add_label(s3, 6.95, 2.27, 5.9, 0.3, "검색 실패 시 ▼", size=12, bold=True, color=ORANGE)

add_box(s3, 0.5, 2.6, 12.3, 1.15, "퍼지 매칭 폴백", [
    "표기 길이 유사 부분문자열 전수비교 (difflib.SequenceMatcher)",
    "개체별 최고점수가 문턱 이상 → 채택 (복수 개체 가능)",
], LORANGE, ORANGE, font_pt=13, title_pt=15.5)

add_arrow(s3, 6.65, 3.75, 6.65, 4.05)

add_box(s3, 0.5, 4.10, 5.95, 2.15, "채택 이유 — 임베딩 대신 문자 매칭", [
    "질의 임베딩 vs 개체 centroid 시도",
    "→ 단일 코퍼스라 centroid 변별력 無",
    "(정답·오답 코사인 차 0.001~0.1)",
    "결론: 폐기",
], LPURPLE, PURPLE, font_pt=12.5, title_pt=14.5)

add_box(s3, 6.85, 4.10, 5.95, 2.15, "0.65 설정 근거 — 실측값", [
    "개체 언급 O: 0.800~1.000",
    "개체 언급 X: 0.333~0.400",
    "0.6 = 동률 구간(오답 기관 혼입)",
    "→ 0.65로 회피",
], LGREEN, GREEN, font_pt=12.5, title_pt=14.5)

footer(s3, "graph_rag/query.py:28 FUZZY_MATCH_THRESHOLD=0.65 · 완전일치 성공 시 퍼지 매칭 생략(중복 계산 방지)")

# ================================================================= Slide 4: 임베더 재사용
s4 = prs.slides.add_slide(BLANK)
add_title(s4, "4. 임베딩/리랭커 모델 메모리 탑재·재사용",
          "STM·MTM·LTM·그래프·문서 검색 — 모델 하나로 통일된 임베딩 공간")

add_box(s4, 0.5, 1.35, 6.0, 3.0, "호출부 5곳 — 동일 factory 경유", [
    "FactMemory — STM/MTM/LTM 검색·승급",
    "DocRetriever — 문서 청크 검색",
    "GraphRetriever.query — 엔티티 퍼지매칭",
    "GraphRetriever.build — 그래프 구축",
    "AgentRuntime — 답변 기여도",
], LBLUE, BLUE, font_pt=13, title_pt=16)

add_arrow(s4, 6.6, 2.85, 7.1, 2.85)

add_box(s4, 7.15, 1.9, 5.7, 1.9, "embedding_factory.build_embedder(cfg)", [
    "캐시 키 = backend·model_id·device·",
    "batch_size·max_seq_length",
    "동일 key → 인스턴스 재사용",
], LGREEN, GREEN, font_pt=13, title_pt=14.5)

add_arrow(s4, 10.0, 3.85, 10.0, 4.15)
add_box(s4, 7.15, 4.20, 5.7, 1.15, "결과: 통일된 임베딩 공간", [
    "5개 저장소 = 동일 벡터공간 → 비교 일관성",
], LPURPLE, PURPLE, font_pt=12.5, title_pt=14)

add_box(s4, 0.5, 4.6, 6.0, 2.35, "실측 낭비 — 캐시 미적용 시", [
    "VRAM 8.47GB",
    "중복분 4.24GB",
    "로딩 중복 7.2초",
    "(arctic-ko 3벌 중복 로드)",
], LRED, RED, font_pt=13.5, title_pt=15.5)

footer(s4, "인스턴스 = 무상태(모델 가중치+캐시 핸들) → 스레드 간 공유 가능 · doc_rag/embedding_factory.py")

# ================================================================= Slide 5: 계층 분기 비교 (카드형)
s5 = prs.slides.add_slide(BLANK)
add_title(s5, "5. STM/MTM/LTM 분기 구조 — 타 아키텍처 비교",
          "현재: 매 턴 세 계층 전부 조회(예산 3/3/3, LLM 판단 아님)")

cards = [
    ("Self-RAG", "판단주체: 모델 자체(파인튜닝)",
     "[Retrieve]/[IsSupported] 토큰 직접 생성", "불가", RED, LRED),
    ("FLARE", "판단주체: 생성 중 신뢰도",
     "토큰 확률 임계치 이하 → 검색 트리거", "불가", RED, LRED),
    ("Adaptive-RAG", "판단주체: 별도 학습 분류기",
     "질의 복잡도 → 검색없음/1단계/다단계", "가능*", ORANGE, LORANGE),
    ("Probing-RAG", "판단주체: 은닉 상태 프로브",
     "은닉층 활성값 → 검색필요 예측 (57.5%↓)", "불가", RED, LRED),
    ("Letta/MemGPT", "판단주체: 토큰 압박 신호",
     "컨텍스트 70%↑ → 경고 → LLM 함수호출", "부분*", ORANGE, LORANGE),
    ("Generative Agents", "판단주체: 분기 없음",
     "recency·importance·relevance 단일 랭킹", "참고", GRAY, BANNER),
]
cw, ch = 3.95, 1.75
gx, gy = 0.25, 0.25
x0, y0 = 0.5, 1.35
for i, (name, who, why, verdict, c, lc) in enumerate(cards):
    col, row = i % 3, i // 3
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    add_box(s5, x, y, cw, ch, name, [who, why], lc, c, font_pt=11.5, title_pt=14.5)
    tag = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + cw - 1.15), Inches(y + 0.08),
                               Inches(1.0), Inches(0.4))
    tag.fill.solid()
    tag.fill.fore_color.rgb = c
    tag.line.fill.background()
    tag.shadow.inherit = False
    tp = tag.text_frame.paragraphs[0]
    tp.text = verdict
    tp.font.size = Pt(12)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER
    tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

concl_y = y0 + 2 * (ch + gy) + 0.05
add_box(s5, 0.5, concl_y, 12.3, 0.85,
        "결론 — 즉시 적용 가능한 방향",
        ["내부접근 필요 방식 배제 → 코사인 하한선(예 0.5) 적용, 미달 계층 제외 (비용 0)"],
        LGREEN, GREEN, font_pt=13, title_pt=15)
add_label(s5, 0.5, concl_y + 0.9, 12.3, 0.3, "* 가능/부분 = 별도 학습·구현 비용 필요", size=11, color=GRAY)

# ================================================================= 부록: 설정값 요약
s6 = prs.slides.add_slide(BLANK)
add_title(s6, "부록 — 설정값 요약", "본문 슬라이드에서 정성적 라벨로 대체된 실제 문턱값")

vals = [
    ("코사인 유사도 문턱", "0.65", "주제 유사 판정 (그래프 퍼지매칭 · LTM 대조 공용)", BLUE, LBLUE),
    ("상호함의 상한", "≥ 2.0", "\"같은 사실(중복)\" 판정", RED, LRED),
    ("상호함의 하한", "< -6.0", "\"표면만 겹친 무관\" 판정 (UPDATE_FLOOR)", ORANGE, LORANGE),
    ("상호함의 갱신 구간", "[-6.0, 2.0)", "\"같은 사안·다른 정보\" 판정 → 기존 LTM 갱신", GREEN, LGREEN),
    ("문서 파일 크기 상한", "120 MB", "초과 시 인제스천 스킵", PURPLE, LPURPLE),
    ("계층별 조회 예산", "3 / 3 / 3", "STM · MTM · LTM 균등 조회 (config.memory_tier_budget)", GRAY, BANNER),
]
vy = 1.5
vh = 0.82
for i, (name, val, note, c, lc) in enumerate(vals):
    add_box(s6, 0.5, vy, 3.3, vh, name, [], lc, c, font_pt=13, title_pt=13.5)
    add_box(s6, 3.95, vy, 2.2, vh, val, [], c, WHITE, font_pt=15, title_pt=17, align=PP_ALIGN.CENTER)
    add_box(s6, 6.3, vy, 6.5, vh, note, [], BANNER, GRAY, font_pt=12.5, title_pt=12.5)
    vy += vh + 0.13

out = r"D:\LAB_RAG\Org-AI-Body\.scratch\design_review.pptx"
prs.save(out)
print("saved:", out)
